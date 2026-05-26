"""theme_extractor 서비스 단위 테스트."""
import os
import io

import pytest

os.environ.setdefault("LLM_PROVIDER", "mock")


def _make_minimal_pptx() -> bytes:
    """테스트용 최소 .pptx 파일 생성."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # 배경 흰색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 제목 텍스트 파란색
    title = slide.shapes.title
    if title:
        title.text = "테스트 제목"
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.mark.asyncio
async def test_extract_from_pptx_returns_theme():
    from app.services.theme_extractor import extract_theme_from_pptx

    pptx_bytes = _make_minimal_pptx()
    theme = await extract_theme_from_pptx(pptx_bytes)

    assert "background" in theme
    assert "primary" in theme
    assert "fontFamily" in theme
    assert theme["source"] == "upload"


@pytest.mark.asyncio
async def test_extract_fallback_on_corrupt_pptx():
    from app.services.theme_extractor import extract_theme_from_pptx

    theme = await extract_theme_from_pptx(b"not a valid pptx")

    # 실패 시 minimal 테마 폴백 반환
    assert theme["background"] == "#fafafa"
    assert theme["source"] == "upload"


@pytest.mark.asyncio
async def test_extract_from_image_returns_fallback_when_provider_returns_invalid(mocker):
    """이미지 Vision 분석 실패(mock provider 비JSON 반환) 시 fallback 반환."""
    # mock provider가 비JSON을 반환하므로 fallback이 반환되어야 한다
    from app.services.theme_extractor import extract_theme_from_image

    result = await extract_theme_from_image(b"fake-image-bytes", "image/png")

    assert "background" in result
    assert "primary" in result
    assert "source" in result
    assert result["source"] == "upload"
