"""테마 추출 서비스 — .pptx 파일 및 이미지에서 색상/폰트 스타일을 추출."""
from __future__ import annotations

import io
import json
from typing import Any

from app.core.logging import get_logger

log = get_logger(__name__)

_MINIMAL_THEME: dict[str, Any] = {
    "source": "upload",
    "background": "#fafafa",
    "primary": "#111827",
    "accent": "#f59e0b",
    "fontFamily": "맑은 고딕",
    "headingSize": 28,
    "bodySize": 14,
    "shapes": {"borderRadius": 0, "strokeColor": "#e5e7eb"},
}


def _rgb_to_hex(rgb: Any) -> str:
    """pptx RGBColor → '#rrggbb'"""
    try:
        return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
    except Exception:
        return "#000000"


async def extract_theme_from_pptx(pptx_bytes: bytes) -> dict[str, Any]:
    """python-pptx로 .pptx 파일에서 배경/주요 색상/폰트를 추출한다."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor

        prs = Presentation(io.BytesIO(pptx_bytes))
        slide = prs.slides[0] if prs.slides else None

        background = "#ffffff"
        primary = "#000000"
        accent = "#2563eb"
        font_family = "맑은 고딕"
        heading_size = 28
        body_size = 14

        if slide:
            # 배경색 추출
            try:
                bg_fill = slide.background.fill
                if bg_fill.type is not None:
                    background = _rgb_to_hex(bg_fill.fore_color.rgb)
            except Exception:
                pass

            # 텍스트 색상 / 폰트 추출 (첫 번째 텍스트 shape 기준)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            primary = _rgb_to_hex(run.font.color.rgb)
                        except Exception:
                            pass
                        if run.font.name:
                            font_family = run.font.name
                        if run.font.size:
                            pt = int(run.font.size / 12700)  # EMU → pt
                            if 8 <= pt <= 72:  # 유효 범위만 사용
                                heading_size = pt
                        break
                    break
                break

        return {
            "source": "upload",
            "background": background,
            "primary": primary,
            "accent": accent,
            "fontFamily": font_family,
            "headingSize": heading_size,
            "bodySize": body_size,
            "shapes": {"borderRadius": 0, "strokeColor": primary},
        }
    except Exception as e:
        log.warning("pptx 스타일 추출 실패, fallback 반환", error=str(e))
        return dict(_MINIMAL_THEME)


async def extract_theme_from_image(image_bytes: bytes, mime: str) -> dict[str, Any]:
    """이미지에서 색상 팔레트/스타일을 분석한다.

    현재 LLM provider가 텍스트 전용 ChatMessage를 사용하므로, Vision API 호출이
    실패하면 _MINIMAL_THEME 폴백을 반환한다. Claude Vision multimodal 지원 시
    ChatMessage에 structured content 블록을 추가해 개선 가능.
    """
    import base64

    from app.providers.llm import ChatMessage, get_llm_provider

    try:
        provider = get_llm_provider()
        b64 = base64.b64encode(image_bytes).decode()
        prompt = (
            "이 이미지의 디자인 스타일을 분석해서 아래 JSON을 반환하십시오.\n"
            "배경색, 주 텍스트 색상, 강조색, 폰트 느낌(한국어 폰트명)을 추출하십시오.\n"
            '{"background":"#hex","primary":"#hex","accent":"#hex","fontFamily":"폰트명","headingSize":28,"bodySize":14}\n'
            "JSON만 반환. 마크다운 없이."
        )
        messages = [
            ChatMessage(
                role="user",
                content=f"data:{mime};base64,{b64}\n\n{prompt}",
            )
        ]
        raw = await provider.complete(messages, temperature=0.1, max_tokens=256)
        parsed = json.loads(raw)
        parsed["source"] = "upload"
        parsed.setdefault("shapes", {"borderRadius": 4, "strokeColor": parsed.get("accent", "#2563eb")})
        return parsed
    except Exception as e:
        log.warning("이미지 Vision 분석 실패, fallback 반환", error=str(e))
        return dict(_MINIMAL_THEME)
